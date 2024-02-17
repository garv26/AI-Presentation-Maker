import streamlit as st
import time
import openai
from pptx import Presentation

# Configure your OpenAI API key
api_key = "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"
openai.api_key = api_key

# Function to generate the presentation
def generate_presentation(topic, num_slides):
    if num_slides > 10:
        num_slides = 10  # Limit the number of slides to a maximum of 5

    prs = Presentation()

    for _ in range(num_slides):
        prompt_subtopic = f"Generate a subtopic related to {topic}."
        response_subtopic = openai.Completion.create(
            model="gpt-3.5-turbo-instruct",
            prompt=prompt_subtopic,
            max_tokens=30,
            n=1
        )
        subtopic = response_subtopic.choices[0].text.strip()

        prompt_content = f"Write about {subtopic} in the context of {topic} in not more than 50 words or in bullet points when needed."
        response_content = openai.Completion.create(
            model="gpt-3.5-turbo-instruct",
            prompt=prompt_content,
            max_tokens=70
        )
        generated_content = response_content.choices[0].text.strip()

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = subtopic
        content = slide.placeholders[1]
        content.text = generated_content

        # Introduce a delay between API requests to stay within the rate limit
        time.sleep(10)  # Adjust the sleep duration based on the rate limit

    prs.save(f"{topic}_presentation.pptx")

# Streamlit app
def main():
    st.title("OpenAI Presentation Generator")

    # User input for the topic and number of slides
    user_topic = st.text_input("Enter the topic for the presentation:")
    user_num_slides = st.number_input("Enter the number of slides:", min_value=1, max_value=10, step=1)

    # Generate the presentation based on user input
    if st.button("Generate Presentation"):
        st.info("Generating presentation... This may take a while.")
        generate_presentation(user_topic, int(user_num_slides))
        st.success(f"Presentation '{user_topic}_presentation.pptx' generated successfully!")

# Run the app
if __name__ == "__main__":
     main()